from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PPTXParser:
    def __init__(self, pptx_path: str):
        """변수 및 관련 모듈 초기화"""
        self.pptx_path = pptx_path
        self.pptx = Presentation(pptx_path)
        self.file_name = Path(pptx_path).name

    def extract_shape(self, shape):
        """shape에서 타입에 맞는 텍스트 추출"""
        if shape.has_text_frame:
            text = shape.text.strip()
            if not text:
                return None
            return {
                "top": shape.top,
                "left": shape.left,
                "content_type": "text",
                "content": shape.text.strip(),
            }
        elif shape.has_table:
            texts = [
                cell.text.strip()
                for row in shape.table.rows
                for cell in row.cells
                if cell.text.strip()
            ]
            return {
                "top": shape.top,
                "left": shape.left,
                "content_type": "text",
                "content": "\n".join(texts),
            }
        elif shape.has_chart and shape.chart.chart_title:
            title = shape.chart.chart_title.text_frame.text.strip()
            return {
                "top": shape.top,
                "left": shape.left,
                "content_type": "text",
                "content": title,
            }
        elif hasattr(shape, "smartArt"):
            texts = [
                node.text_frame.text.strip()
                for node in shape.smartArt.nodes
                if node.text_frame.text.strip()
            ]
            return {
                "top": shape.top,
                "left": shape.left,
                "content_type": "text",
                "content": "\n".join(texts),
            }
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_items = []
            for sub_shape in shape.shapes:
                res = self.extract_shape(sub_shape)
                if res:
                    group_items.append(res)
            group_items.sort(key=lambda x: (x["top"], x["left"]))
            return {
                "top": shape.top,
                "left": shape.left,
                "content_type": "group",
                "content": group_items,
            }
        return None

    def parse_slides(self):
        """파일에서 슬라이드 단위 텍스트 파싱"""
        all_slides = []
        for slide_idx, slide in enumerate(self.pptx.slides):
            items = []
            for shape in slide.shapes:
                extract_result = self.extract_shape(shape)
                if extract_result:
                    items.append(extract_result)
            items.sort(key=lambda x: (x["top"], x["left"]))
            slide_info = {
                "file_name": self.file_name,
                "file_path": self.pptx_path,
                "index": slide_idx,
                "text_content": self._extract_text_from_slide(items),
            }
            all_slides.append(slide_info)
        return all_slides

    def _extract_text_from_slide(self, items):
        """pptx-parser 모듈 통해 슬라이드 단위로 추출된 데이터를 입력받아 연결된 텍스트로 변환"""

        def extract_text_recursive(items):
            texts = []
            for item in items:
                if item["content_type"] == "text":
                    texts.append(item["content"])
                elif item["content_type"] == "group":
                    texts.extend(extract_text_recursive(item["content"]))
            return texts

        return "\n".join(extract_text_recursive(items))
